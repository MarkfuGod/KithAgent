"""
File content extractors for non-plaintext formats.

Handles Word (.docx), PDF, and images. Each extractor returns
a dict with 'type' (text|image) and the extracted content so
the summarizer knows whether to use a text or vision model.
"""

from __future__ import annotations

import base64
import logging
from pathlib import Path

logger = logging.getLogger("agent_sys.extractors")

_IMAGE_EXTENSIONS = frozenset({
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp",
})

_DOCUMENT_EXTENSIONS = frozenset({
    ".pdf", ".docx", ".doc", ".pptx", ".xlsx",
})

_TEXT_EXTENSIONS = frozenset({
    ".py", ".js", ".ts", ".md", ".txt", ".json", ".yaml", ".yml",
    ".toml", ".sh", ".go", ".rs", ".html", ".css", ".jsx", ".tsx",
    ".java", ".c", ".cpp", ".h", ".hpp", ".rb", ".php", ".swift",
    ".kt", ".scala", ".r", ".sql", ".xml", ".csv", ".ini", ".cfg",
    ".env", ".log", ".rst", ".tex", ".vue", ".svelte",
})


def is_image(ext: str) -> bool:
    return ext.lower() in _IMAGE_EXTENSIONS


def is_document(ext: str) -> bool:
    return ext.lower() in _DOCUMENT_EXTENSIONS


def is_plaintext(ext: str) -> bool:
    return ext.lower() in _TEXT_EXTENSIONS


def extract_text_from_pdf(path: str, max_chars: int = 4000) -> str | None:
    """Extract text content from a PDF file."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
            if sum(len(t) for t in text_parts) >= max_chars:
                break
        doc.close()
        text = "\n".join(text_parts)
        return text[:max_chars] if text.strip() else None
    except ImportError:
        pass

    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(path)
        text_parts = []
        for page in reader.pages:
            text_parts.append(page.extract_text() or "")
            if sum(len(t) for t in text_parts) >= max_chars:
                break
        text = "\n".join(text_parts)
        return text[:max_chars] if text.strip() else None
    except ImportError:
        logger.debug("No PDF library available (install PyMuPDF or PyPDF2)")
        return None
    except Exception as e:
        logger.debug("PDF extraction failed for %s: %s", path, e)
        return None


def extract_text_from_docx(path: str, max_chars: int = 4000) -> str | None:
    """Extract text content from a Word .docx file."""
    try:
        from docx import Document
        doc = Document(path)
        text_parts = []
        total = 0
        for para in doc.paragraphs:
            text_parts.append(para.text)
            total += len(para.text)
            if total >= max_chars:
                break
        text = "\n".join(text_parts)
        return text[:max_chars] if text.strip() else None
    except ImportError:
        logger.debug("python-docx not installed — cannot extract .docx")
        return None
    except Exception as e:
        logger.debug("DOCX extraction failed for %s: %s", path, e)
        return None


def extract_text_from_pptx(path: str, max_chars: int = 4000) -> str | None:
    """Extract text content from a PowerPoint .pptx file (best-effort)."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text_parts: list[str] = []
        total = 0
        for slide_idx, slide in enumerate(prs.slides, 1):
            text_parts.append(f"[slide {slide_idx}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
                    total += len(shape.text)
            if total >= max_chars:
                break
        text = "\n".join(text_parts)
        return text[:max_chars] if text.strip() else None
    except ImportError:
        logger.debug("python-pptx not installed — cannot extract .pptx")
        return None
    except Exception as e:
        logger.debug("PPTX extraction failed for %s: %s", path, e)
        return None


def extract_text_from_xlsx(path: str, max_chars: int = 4000, max_rows_per_sheet: int = 50) -> str | None:
    """Extract text content from an Excel .xlsx file as header + sample rows.

    We don't want to serialize a 100k-row sheet, so we just sample.
    """
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        text_parts: list[str] = []
        total = 0
        for sheet in wb.worksheets:
            text_parts.append(f"[sheet: {sheet.title}]")
            for i, row in enumerate(sheet.iter_rows(values_only=True)):
                if i >= max_rows_per_sheet:
                    text_parts.append(f"... ({sheet.max_row} total rows)")
                    break
                line = " | ".join(str(c) if c is not None else "" for c in row)
                text_parts.append(line)
                total += len(line)
                if total >= max_chars:
                    break
            if total >= max_chars:
                break
        wb.close()
        text = "\n".join(text_parts)
        return text[:max_chars] if text.strip() else None
    except ImportError:
        logger.debug("openpyxl not installed — cannot extract .xlsx")
        return None
    except Exception as e:
        logger.debug("XLSX extraction failed for %s: %s", path, e)
        return None


def encode_image_base64(path: str, max_bytes: int = 5 * 1024 * 1024) -> str | None:
    """Read an image file and return its base64-encoded data URI."""
    try:
        p = Path(path)
        if not p.exists() or p.stat().st_size > max_bytes:
            return None

        raw = p.read_bytes()
        b64 = base64.b64encode(raw).decode("ascii")

        ext = p.suffix.lower()
        mime_map = {
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".webp": "image/webp",
            ".bmp": "image/bmp",
        }
        mime = mime_map.get(ext, "image/png")
        return f"data:{mime};base64,{b64}"
    except Exception as e:
        logger.debug("Image encoding failed for %s: %s", path, e)
        return None


def render_pdf_page_base64(
    path: str,
    page_number: int = 0,
    max_bytes: int = 5 * 1024 * 1024,
) -> str | None:
    """Render one PDF page to a PNG data URI for vision/multimodal models."""
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(path)
        try:
            if page_number < 0 or page_number >= len(doc):
                return None
            page = doc[page_number]
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            raw = pix.tobytes("png")
        finally:
            doc.close()
        if len(raw) > max_bytes:
            return None
        b64 = base64.b64encode(raw).decode("ascii")
        return f"data:image/png;base64,{b64}"
    except ImportError:
        logger.debug("PyMuPDF not installed — cannot render PDF page image")
        return None
    except Exception as e:
        logger.debug("PDF page rendering failed for %s: %s", path, e)
        return None


def extract_content(path: str, max_chars: int = 4000) -> dict | None:
    """
    Extract content from any supported file type.

    Returns:
        {"type": "text", "content": "..."} for text-extractable files
        {"type": "image", "data_uri": "data:image/png;base64,..."} for images
        None if extraction fails
    """
    ext = Path(path).suffix.lower()

    if ext == ".pdf":
        text = extract_text_from_pdf(path, max_chars)
        if text:
            return {"type": "text", "content": text}
        # Image-only PDFs need a rendered page, not raw PDF bytes.
        data_uri = render_pdf_page_base64(path)
        if data_uri:
            return {"type": "image", "data_uri": data_uri, "metadata": {"source_kind": "scanned_pdf", "page": 1}}
        return None

    if ext in (".docx", ".doc"):
        if ext == ".docx":
            text = extract_text_from_docx(path, max_chars)
            if text:
                return {"type": "text", "content": text}
        # .doc (old binary Word) isn't supported — skip cleanly.
        return None

    if ext == ".pptx":
        text = extract_text_from_pptx(path, max_chars)
        if text:
            return {"type": "text", "content": text}
        return None

    if ext == ".xlsx":
        text = extract_text_from_xlsx(path, max_chars)
        if text:
            return {"type": "text", "content": text}
        return None

    if is_image(ext):
        data_uri = encode_image_base64(path)
        if data_uri:
            return {"type": "image", "data_uri": data_uri}
        return None

    if is_plaintext(ext):
        try:
            text = Path(path).read_text(errors="replace")
            return {"type": "text", "content": text[:max_chars]}
        except Exception:
            return None

    return None
